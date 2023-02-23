import pprint
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, random

img_file_names = []
quizes = []

for filename in os.listdir("four-letter-game/data/imgs"):
    if filename not in img_file_names:
        img_file_names.append(filename)

def read_data():
    global quizes
    filename = "four-letter-game/data/data.txt"
    f = open(filename, "r", encoding="utf-8")

    lines = f.readlines()
    for line in lines:
        quizes.append(line)

def ppt_gen():
    prs = Presentation()
    size = len(quizes)+len(img_file_names)
    for i in range(size):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        tmp = random.randint(6,10)
        if i%tmp == 0:
            if not len(img_file_names) == 0:
                img_path = random.choice(img_file_names)
                pic = slide.shapes.add_picture("four-letter-game/data/imgs/"+img_path, Inches(1), Inches(1), height=Inches(5))
                pic.left = int((prs.slide_width - pic.width) / 2)
                pic.top = int((prs.slide_height - pic.height) / 2)
                img_file_names.remove(img_path)
            else:
                # add 4글자 문제
                txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
                tf = txBox.text_frame
                randomQuiz = random.choice(quizes)
                tf.text = randomQuiz[0:2]
                
                tf.paragraphs[0].font.size = Pt(100)
                tf.paragraphs[0].font.bold = True

                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                txBox.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                quizes.remove(randomQuiz)  
        else:
            # add 4글자 문제
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
            tf = txBox.text_frame
            randomQuiz = random.choice(quizes)
            tf.text = randomQuiz[0:2]
            
            tf.paragraphs[0].font.size = Pt(100)
            tf.paragraphs[0].font.bold = True

            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            txBox.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            quizes.remove(randomQuiz)
    
    prs.save("4글자 이어말하기.pptx")

read_data()
ppt_gen()