import random

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, random

def ppt_gen():
    prs = Presentation()
    for num in range(30):
        for char in ["A", "B", "C", "D", "E"]:
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)

            # add 4글자 문제
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
            tf = txBox.text_frame
            tf.text = f'{num}{char}'
            
            tf.paragraphs[0].font.size = Pt(300)
            tf.paragraphs[0].font.bold = True

            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            txBox.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
    prs.save("로또번호.pptx")

ppt_gen()






