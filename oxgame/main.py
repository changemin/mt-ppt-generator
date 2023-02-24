import pprint
from pptx import Presentation

quizes = []

def read_data():
    global quizes
    filename = "oxgame/data.txt"
    f = open(filename, "r", encoding="utf-8")


    lines = f.readlines()
    for line in lines:
        quiz = {
            "type": line[1],
            "quiz": line[3:len(line)-2],
            "ans": line[len(line)-2]
        }
        quizes.append(quiz)
    #     print("Type: ", quiz["type"], "Quiz: ", quiz["quiz"], "Ans: ", quiz["ans"])
    # print("@"*20)
# pprint.pprint(quizes)
    
def ppt_gen():
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # add some text to the slide
    title = slide.shapes.title
    title.text = "My Title"
    subtitle = slide.placeholders[1]
    subtitle.text = "My Subtitle"

    # save the presentation to a file
    prs.save("output.pptx")

read_data()
print(quizes)

f = open("oxgame/data2.csv", "w", encoding="utf-8")
msg = ""
for quiz in quizes:
    msg += f'{quiz["ans"]}\n'
print(msg)
f.write(msg)
f.close()

# ppt_gen()