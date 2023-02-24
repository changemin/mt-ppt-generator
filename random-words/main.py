import words
import random

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, random

myWords = []

for word in words.data["result"]:
    myWords.append(word["name"])
    
random.shuffle(myWords)

print(len(words.data["result"]))

def ppt_gen():
    prs = Presentation()
    for word in list(myWords):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

       
        # add 4글자 문제
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
        tf = txBox.text_frame
        tf.text = word
        
        tf.paragraphs[0].font.size = Pt(100)
        tf.paragraphs[0].font.bold = True

        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        txBox.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        myWords.remove(word)  
        
    
    prs.save("제시어뭉탱이.pptx")

ppt_gen()






